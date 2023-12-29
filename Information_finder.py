#!/usr/bin/env python
# coding: utf-8

# In[11]:


import os
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import io
from docx import Document
from pptx import Presentation
import concurrent.futures

class DocumentProcessor:
    def __init__(self, search_terms, similarity_threshold=0.9, pdf_iterations=2):
        self.search_terms = search_terms
        self.similarity_threshold = similarity_threshold
        self.pdf_iterations = pdf_iterations

    def text_similarity(self, text1, text2):
        matches = sum(c1 == c2 for c1, c2 in zip(text1, text2))
        return matches / max(len(text1), len(text2))

    def extract_paragraphs(self, text):
        return text.split('\n\n')

    def find_term_in_paragraphs(self, paragraphs, term):
        lower_term = term.lower()
        for paragraph in paragraphs:
            if lower_term in paragraph.lower():
                return paragraph
        return None

    def process_pdf(self, file_path):
        best_text = ""
        best_similarity = 0
        base_filename = os.path.splitext(file_path)[0]

        for iteration in range(self.pdf_iterations):
            doc = fitz.open(file_path)
            standard_text = ''.join([page.get_text() for page in doc])

            ocr_text = ''
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap()
                img = Image.open(io.BytesIO(pix.tobytes()))
                ocr_text += pytesseract.image_to_string(img)

            similarity = self.text_similarity(standard_text, ocr_text)
            selected_text = standard_text if similarity >= self.similarity_threshold else ocr_text

            if similarity > best_similarity:
                best_text = selected_text
                best_similarity = similarity

            # Save the iteration text to a temporary file
            with open(f"{base_filename}_iteration{iteration}.txt", "w") as file:
                file.write(selected_text)

        # Keep the best text and delete other iterations
        with open(f"{base_filename}.txt", "w") as file:
            file.write(best_text)

        for iteration in range(self.pdf_iterations):
            os.remove(f"{base_filename}_iteration{iteration}.txt")

        return f"{base_filename}.txt"

    def process_docx(self, file_path):
        doc = Document(file_path)
        text = '\n'.join([para.text for para in doc.paragraphs])
        with open(f"{os.path.splitext(file_path)[0]}.txt", "w") as file:
            file.write(text)
        return f"{os.path.splitext(file_path)[0]}.txt"

    def process_pptx(self, file_path):
        ppt = Presentation(file_path)
        text = '\n'.join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")])
        with open(f"{os.path.splitext(file_path)[0]}.txt", "w") as file:
            file.write(text)
        return f"{os.path.splitext(file_path)[0]}.txt"

    def search_terms_in_text(self, file_path):
        results = []
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            text = file.read()
            paragraphs = self.extract_paragraphs(text)

            for term in self.search_terms:
                found_paragraph = self.find_term_in_paragraphs(paragraphs, term)
                if found_paragraph:
                    results.append(f"Found '{term}' in {file_path}:\n{found_paragraph}\n\n")
                else:
                    results.append(f"'{term}' not found in {file_path} - marked as impossible.\n\n")
        
        return results

def combine_text_files(folder_path, output_filename="ALL.txt"):
    with open(os.path.join(folder_path, output_filename), "w", encoding='utf-8', errors='ignore') as outfile:
        for filename in os.listdir(folder_path):
            if filename.endswith('.txt') and filename != output_filename:
                try:
                    with open(os.path.join(folder_path, filename), 'r', encoding='utf-8') as infile:
                        outfile.write(infile.read() + "\n")
                except UnicodeDecodeError:
                    with open(os.path.join(folder_path, filename), 'r', encoding='ISO-8859-1') as infile:
                        outfile.write(infile.read() + "\n")

def process_file(processor, folder_path, filename):
    file_path = os.path.join(folder_path, filename)
    if filename.endswith('.pdf'):
        return processor.process_pdf(file_path)
    elif filename.endswith('.docx'):
        return processor.process_docx(file_path)
    elif filename.endswith('.pptx'):
        return processor.process_pptx(file_path)
    return None

def main():
    search_terms = [
    ]

    processor = DocumentProcessor(search_terms)
    folder_path = os.getcwd()
    all_results = []

    with concurrent.futures.ThreadPoolExecutor() as executor:
        futures = {executor.submit(process_file, processor, folder_path, filename): filename for filename in os.listdir(folder_path)}
        for future in concurrent.futures.as_completed(futures):
            text_file = future.result()
            if text_file:
                results = processor.search_terms_in_text(text_file)
                all_results.extend(results)

    with open(os.path.join(folder_path, "search_results.txt"), "w", encoding='utf-8', errors='ignore') as result_file:
        result_file.write("\n".join(all_results))

    combine_text_files(folder_path)

if __name__ == "__main__":
    main()


# In[ ]:


search_terms = [
    "HOLC maps","HOLC", "Tompkins Square riot of 1988", "Festival markets", "Megalopolis",
    "Hope VI housing policy","Hope VI", "Artificial Levees in New Orleans", "Containerization in shipping",
    "US Housing Act of 1949", "Levittown", "Gentrification", "Demolition of Pruitt-Igoe",
    "Sprawl", "Annexation of suburban territory", "Restrictive covenants", "Demolition of Pennsylvania Station, New York",
    "Riots of the 1960s", "Ford to City: Drop Dead", "Deindustrialization", "Peak of violent crime in US cities",
    "Block busting", "Great Migration", "Second ghetto", "Rust Belt", "Edge Cities",
    "Superstorm Sandy", "Growth of the Sun Belt"
    ]

